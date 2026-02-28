from __future__ import annotations

import io
import re
import unicodedata
from typing import Any, Dict, List, Tuple, Optional

from fastapi import UploadFile

from app.core.config import settings
from app.services.text_repair import repair_ocr_spacing_line, repair_ocr_spacing_text, fix_eth_d


# Postgres TEXT/VARCHAR cannot contain NUL (\x00) bytes.
# Some extractors (notably PDF text extraction) may produce NUL/control chars.
_CTRL_RE = re.compile(r"[\x00-\x08\x0B\x0C\x0E-\x1F\x7F]")


def _sanitize_text(text: str) -> str:
    """Remove characters that can break DB storage / downstream processing."""
    if not text:
        return ""
    # Replace control chars (including NUL) with spaces, keep \t \n \r.
    text = _CTRL_RE.sub(" ", text)
    return text


def _normalize_pipeline_text(text: str) -> str:
    """Normalize extracted text for stable Vietnamese display/storage."""
    if not text:
        return ""
    text = _sanitize_text(text)
    text = unicodedata.normalize("NFKC", text)
    text = fix_eth_d(text)
    text = repair_ocr_spacing_text(text)
    return _sanitize_text(text)


# Some PDFs insert soft hyphens or glue single-letter variables into words.
_SOFT_HYPHEN = "\u00ad"
_MATH_VARS = "xyzuvwtnmk"


def _fix_joined_variables(text: str) -> str:
    """Best-effort fix for common PDF extraction artifacts.

    Examples we try to repair:
      - "biếnyđược" -> "biến y được"
      - "từxqua" -> "từ x qua"

    This is intentionally conservative and targets single-letter math variables.
    """
    if not text:
        return ""
    s = text.replace(_SOFT_HYPHEN, "")
    # normalize unicode (keeps Vietnamese diacritics stable)
    s = unicodedata.normalize("NFKC", s)
    # word + var + word
    s = re.sub(rf"([A-Za-zÀ-ỹà-ỹ]{{2,}})([{_MATH_VARS}])([A-Za-zÀ-ỹà-ỹ]{{2,}})", r"\1 \2 \3", s)
    # word + var (end)
    s = re.sub(rf"([A-Za-zÀ-ỹà-ỹ]{{2,}})([{_MATH_VARS}])\b", r"\1 \2", s)
    # var + word
    s = re.sub(rf"\b([{_MATH_VARS}])([A-Za-zÀ-ỹà-ỹ]{{2,}})", r"\1 \2", s)
    return s


def _has_vn_diacritics(s: str) -> bool:
    # quick check: any non-ascii letter (covers Vietnamese diacritics)
    return any(ord(ch) > 127 for ch in (s or ""))


def _repair_spaced_letters_line(line: str) -> str:
    """Backward-compatible wrapper.

    We now use a stronger repair that also fixes short split patterns like:
      - "Lập t rình" -> "Lập trình"
      - "điề u k hiể n" -> "điều khiển"
    """
    return repair_ocr_spacing_line(line)




# ---------------------
# PDF cleaning helpers
# ---------------------
_PDF_FOOTER_PATTERNS = [
    # e.g. "Lập trình Python dành cho kỹ thuật - Trang 24"
    re.compile(r".*?\bTrang\s+\d+\b.*$", flags=re.IGNORECASE),
    # standalone page number
    re.compile(r"^\s*\d+\s*$"),
]
_TOC_HINT_RE = re.compile(r"(m\s*ục\s*l\s*ục|table\s*of\s*contents)", flags=re.IGNORECASE)
_DOTS_LEADER_RE = re.compile(r"\.{3,}")  # .........


def _clean_pdf_page_text(text: str) -> str:
    """Remove common PDF headers/footers and noisy lines."""
    if not text:
        return ""
    text = _sanitize_text(text)
    # Normalize line endings
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    lines = []
    for raw in text.split("\n"):
        line = raw.strip()
        if not line:
            continue
        line = _fix_joined_variables(line)
        line = repair_ocr_spacing_line(line)
        # Drop leader-dot lines common in TOC
        if _DOTS_LEADER_RE.search(line) and len(line) < 120:
            continue
        # Drop headers/footers
        if any(pat.match(line) for pat in _PDF_FOOTER_PATTERNS):
            continue
        lines.append(line)
    return _normalize_pipeline_text("\n".join(lines).strip())


def _extract_text_pdf_ocr(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    """OCR fallback for scanned PDFs (image-only).

    Uses PyMuPDF for rendering + Tesseract via pytesseract.
    """
    if not bool(getattr(settings, "PDF_OCR_ENABLED", True)):
        return "", []

    try:
        import fitz  # PyMuPDF
    except Exception:
        return "", []

    try:
        from PIL import Image
    except Exception:
        return "", []

    ocr_backend = None
    pytesseract = None
    easyocr = None
    try:
        import pytesseract  # type: ignore
        ocr_backend = 'tesseract'
    except Exception:
        try:
            import easyocr  # type: ignore
            ocr_backend = 'easyocr'
        except Exception:
            return "", []

    lang = str(getattr(settings, "PDF_OCR_LANG", "vie+eng") or "vie+eng")
    zoom = float(getattr(settings, "PDF_OCR_ZOOM", 2.5) or 2.5)
    max_pages = int(getattr(settings, "PDF_OCR_MAX_PAGES", 200) or 200)

    doc = fitz.open(stream=data, filetype="pdf")
    pages_out: List[Dict[str, Any]] = []
    full_parts: List[str] = []

    reader = None
    if ocr_backend == 'easyocr':
        # EasyOCR uses language codes like ['vi','en']
        langs = []
        low = lang.lower()
        if 'vie' in low or 'vi' in low:
            langs.append('vi')
        if 'eng' in low or 'en' in low:
            langs.append('en')
        if not langs:
            langs = ['en']
        try:
            reader = easyocr.Reader(langs, gpu=False)
        except Exception:
            reader = None

    for idx in range(min(len(doc), max_pages)):
        page = doc.load_page(idx)
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
        if ocr_backend == 'tesseract':
            raw_text = pytesseract.image_to_string(img, lang=lang) or ""
        else:
            if reader is None:
                continue
            raw_text = "\n".join(reader.readtext(img, detail=0, paragraph=True) or [])
        page_text = _clean_pdf_page_text(raw_text)
        if not page_text:
            continue
        if _is_probable_toc(page_text):
            continue
        full_parts.append(page_text)
        for ch in _chunk_text(page_text):
            pages_out.append({"text": ch, "meta": {"page": idx + 1, "ocr": True, "ocr_backend": ocr_backend}})

    full_text = _normalize_pipeline_text("\n\n".join(full_parts))
    return full_text, pages_out


def _is_probable_toc(page_text: str) -> bool:
    """Heuristic to skip Table of Contents pages."""
    if not page_text:
        return False
    low = page_text.lower()
    lines = [ln.strip() for ln in page_text.splitlines() if ln.strip()]
    if not lines:
        return False
    short_lines = sum(1 for ln in lines if len(ln) <= 80) / max(1, len(lines))
    dotted_lines = sum(1 for ln in lines if _DOTS_LEADER_RE.search(ln)) / max(1, len(lines))
    digits = sum(ch.isdigit() for ch in page_text)
    digit_ratio = digits / max(1, len(page_text))
    chapter_hits = low.count("chương") + low.count("chapter")
    has_toc_hint = bool(_TOC_HINT_RE.search(low))

    # Require stronger combined signals to reduce false positives from normal content lines.
    if has_toc_hint and dotted_lines >= 0.18 and short_lines >= 0.55:
        return True
    if has_toc_hint and chapter_hits >= 3 and short_lines >= 0.55:
        return True
    if dotted_lines >= 0.30 and short_lines >= 0.65 and digit_ratio >= 0.08:
        return True
    return False


def _candidate_page_coverage(chunks: List[Dict[str, Any]], total_pages: int | None = None) -> float:
    pages = []
    for ch in chunks or []:
        meta = (ch or {}).get("meta") or {}
        p = meta.get("page")
        try:
            if p is not None:
                pages.append(int(p))
        except Exception:
            continue
    if not pages:
        return 0.0
    unique = len(set(pages))
    max_page = max(pages)
    cov_pages = max(unique, max_page)
    if total_pages and int(total_pages) > 0:
        return min(1.0, cov_pages / float(total_pages))
    return float(cov_pages)

def _chunk_text(text: str, chunk_size: int = 900, overlap: int = 120) -> List[str]:
    """Smart chunking that avoids breaking across topic/section boundaries.

    Why this matters:
      The app stores topics by (start_chunk_index, end_chunk_index). If the chunker slices
      purely by characters, a single chunk can contain multiple major headings (e.g. many
      "Chủ đề:" sections). Topic extraction then merges several topics into one because
      it sees only the first heading of the chunk.

    Strategy (deterministic):
      1) Normalize whitespace but preserve newlines.
      2) If we see strong separator lines (====...), split into segments.
      3) Otherwise, split into segments at strong heading lines.
      4) Chunk each segment with overlap, WITHOUT crossing segment boundaries.
    """

    def _chunk_chars(seg: str) -> List[str]:
        seg = seg.strip()
        if not seg:
            return []
        out: List[str] = []
        start = 0
        while start < len(seg):
            end = min(len(seg), start + chunk_size)
            # Prefer cutting on a newline near the end so we don't split headings/markers.
            if end < len(seg):
                # search a window near the end for a safe break
                win_start = max(start + 200, end - 220)
                cut = seg.rfind("\n", win_start, end)
                if cut != -1 and cut > start + 200:
                    end = cut
            out.append(seg[start:end])
            if end == len(seg):
                break
            start = max(0, end - overlap)
        return out

    text = _sanitize_text(text)
    # Preserve newlines (important for heading detection/topic extraction).
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # Collapse excessive spaces but keep line breaks.
    text = re.sub(r"[\t\f\v ]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = text.strip()
    if not text:
        return []

    # 1) Strong visual separators (common in TXT exports) → hard boundaries.
    # Example: "============================================================"
    sep_rx = re.compile(r"\n\s*[=\-_]{8,}\s*\n")

    # 2) Strong heading lines (VN/EN). Kept conservative to avoid splitting on exercises.
    heading_rx = re.compile(
        r"^\s*(?:[\-•\*]+\s*)?(?:(?:\d{1,2}(?:\.\d{1,3}){0,4}|[IVXLCDM]{1,8})\s*[\)\.\:\-–]\s*)?(?:"
        r"(?:chủ\s*đề|chu\s*de|topic)\s*[:\-–]\s+\S+"
        r"|(?:chương|chuong|chapter|unit|lesson)\s+(?:[0-9]{1,3}|[IVXLCDM]{1,6})\b"
        r"|(?:phần|phan|mục|muc|bài|bai|section)\s+(?:[0-9]{1,3}(?:\.[0-9]{1,3}){0,4}|[IVXLCDM]{1,6})\b"
        r"|(?:phụ\s*lục|phu\s*luc|appendix)\s*(?:[:\-–]|—)\s+\S+"
        r"|(?:giới\s*thiệu|gioi\s*thieu|introduction)\b\s*$"
        r"|(?:phần|phan)\s+[\"“”'`].+[\"“”'`]\s*\(.*\)\s*$"
        r")",
        flags=re.IGNORECASE,
    )

    # "Soft" subject headings (no numeric indices / no "Chủ đề:" label).
    # Used to prevent chunking across subject boundaries in teacher-style TXT docs.
    marker_rx = re.compile(r"^\s*(ít\s*dữ\s*liệu|it\s*du\s*lieu|ý\s*chính|y\s*chinh|khái\s*niệm|khai\s*niem|mục\s*tiêu|muc\s*tieu|quiz-ready)\b", re.IGNORECASE)
    subject_hints = [
        'toán', 'đại số', 'giải tích',
        'vật lý', 'cơ học', 'điện học', 'nhiệt học',
        'hóa', 'hóa học', 'phản ứng', 'dung dịch',
        'sinh', 'sinh học', 'tế bào', 'di truyền', 'sinh thái',
        'tin', 'tin học', 'thuật toán', 'dữ liệu', 'an toàn thông tin',
        'xác suất', 'thống kê',
        'kinh tế', 'cung', 'cầu', 'lạm phát', 'lãi suất',
        'địa lý', 'khí hậu', 'dân số', 'tài nguyên',
        'lịch sử', 'bối cảnh', 'nguyên nhân', 'hệ quả',
        'đọc hiểu', 'viết', 'ngữ văn',
        'kỹ năng', 'học tập', 'ghi nhớ',
    ]

    def _next_non_empty(lines_in: list[str], idx: int, max_ahead: int = 3) -> str | None:
        for j in range(idx + 1, min(len(lines_in), idx + 1 + int(max_ahead))):
            s = (lines_in[j] or '').strip()
            if s:
                return s
        return None

    def _is_soft_heading_line(ln: str, next_ln: str | None) -> bool:
        s = (ln or '').strip()
        if not s:
            return False
        # avoid splitting on sentences / definitions
        if any(x in s for x in ['=', '→', '->', '⇒']):
            return False
        if ':' in s or '：' in s:
            return False
        if re.search(r"[\.!?;]$", s):
            return False
        if len(s) < 8 or len(s) > 95:
            return False
        sl = s.lower()
        has_shape = (('(' in s and ')' in s) or ('–' in s) or ('—' in s) or (' - ' in s) or any(h in sl for h in subject_hints))
        if not has_shape:
            return False
        if not next_ln:
            return False
        return bool(marker_rx.match(next_ln.strip()))

    segments: List[str] = []

    if sep_rx.search(text):
        parts = [p.strip() for p in sep_rx.split(text) if p and p.strip()]
        segments = parts
    else:
        # Split by heading lines (do not drop the heading line; it belongs to the next segment).
        lines = text.split("\n")
        buf: List[str] = []
        for i, ln in enumerate(lines):
            nxt = _next_non_empty(lines, i, 3)
            if heading_rx.match((ln or "").strip()) or _is_soft_heading_line(ln, nxt):
                if buf and "\n".join(buf).strip():
                    segments.append("\n".join(buf).strip())
                buf = [ln]
            else:
                buf.append(ln)
        if buf and "\n".join(buf).strip():
            segments.append("\n".join(buf).strip())

        # Merge tiny segments (often OCR artefacts) into previous to avoid over-splitting.
        merged: List[str] = []
        for seg in segments:
            if not merged:
                merged.append(seg)
                continue
            if len(seg) < 220:
                merged[-1] = (merged[-1] + "\n" + seg).strip()
            else:
                merged.append(seg)
        segments = merged

    chunks: List[str] = []
    for seg in segments:
        chunks.extend(_chunk_chars(seg))

    return [c for c in chunks if c and c.strip()]


def _extract_text_pdf_pypdf(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text using pypdf (fallback)."""
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages_out: List[Dict[str, Any]] = []
    full_parts: List[str] = []

    for idx, page in enumerate(reader.pages):
        raw_text = page.extract_text() or ""
        page_text = _clean_pdf_page_text(raw_text)
        if not page_text:
            continue
        if _is_probable_toc(page_text):
            continue
        full_parts.append(page_text)
        for ch in _chunk_text(page_text):
            pages_out.append({"text": ch, "meta": {"page": idx + 1}})

    full_text = _normalize_pipeline_text("\n\n".join(full_parts))
    return full_text, pages_out


def _extract_text_pdf_pymupdf(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text using PyMuPDF (fitz). Often best for PDF with proper text layer."""
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    pages_out: List[Dict[str, Any]] = []
    full_parts: List[str] = []

    for idx in range(len(doc)):
        page = doc.load_page(idx)
        raw_text = page.get_text("text") or ""
        page_text = _clean_pdf_page_text(raw_text)
        if not page_text:
            continue
        if _is_probable_toc(page_text):
            continue
        full_parts.append(page_text)
        for ch in _chunk_text(page_text):
            pages_out.append({"text": ch, "meta": {"page": idx + 1}})

    full_text = _normalize_pipeline_text("\n\n".join(full_parts))
    return full_text, pages_out


def _extract_text_pdf_pymupdf_words(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text using PyMuPDF 'words' reconstruction.

    This often fixes missing spaces that happen with certain PDFs (especially math/variables).
    """
    import fitz  # PyMuPDF

    doc = fitz.open(stream=data, filetype="pdf")
    pages_out: List[Dict[str, Any]] = []
    full_parts: List[str] = []

    for idx in range(len(doc)):
        page = doc.load_page(idx)
        words = page.get_text("words") or []
        if not words:
            continue
        # word tuple: x0, y0, x1, y1, "word", block_no, line_no, word_no
        words.sort(key=lambda w: (w[5], w[6], w[7], w[0]))
        lines: List[str] = []
        cur_key = None
        cur: List[str] = []
        for w in words:
            key = (w[5], w[6])
            txt = str(w[4] or "").strip()
            if not txt:
                continue
            if cur_key is None:
                cur_key = key
            if key != cur_key:
                if cur:
                    lines.append(" ".join(cur))
                cur = [txt]
                cur_key = key
            else:
                cur.append(txt)
        if cur:
            lines.append(" ".join(cur))

        raw_text = "\n".join(lines)
        page_text = _clean_pdf_page_text(raw_text)
        if not page_text:
            continue
        if _is_probable_toc(page_text):
            continue
        full_parts.append(page_text)
        for ch in _chunk_text(page_text):
            pages_out.append({"text": ch, "meta": {"page": idx + 1}})

    full_text = _normalize_pipeline_text("\n\n".join(full_parts))
    return full_text, pages_out


def _extract_text_pdf_pdfplumber(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text using pdfplumber (pdfminer.six). Useful when other extractors miss text."""
    import logging
    import pdfplumber

    # Silence noisy pdfminer warnings (does not affect extraction quality)
    logging.getLogger("pdfminer").setLevel(logging.ERROR)

    pages_out: List[Dict[str, Any]] = []
    full_parts: List[str] = []

    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for idx, page in enumerate(pdf.pages):
            raw_text = page.extract_text() or ""

            # Sparse text often indicates complex multi-column layouts; try alternate extractors.
            sparse = len((raw_text or "").strip()) < 220
            if sparse:
                candidates = [raw_text or ""]
                for kwargs in (
                    {"use_text_flow": True},
                    {"layout": True},
                    {"x_tolerance": 1.5, "y_tolerance": 2.0},
                    {"use_text_flow": True, "x_tolerance": 1.5, "y_tolerance": 2.0},
                ):
                    try:
                        alt = page.extract_text(**kwargs) or ""
                        if alt:
                            candidates.append(alt)
                    except Exception:
                        continue

                try:
                    from app.services.text_quality import quality_score
                except Exception:
                    quality_score = None  # type: ignore

                best = raw_text or ""
                best_s = -1e9
                for cand in candidates:
                    cleaned = _clean_pdf_page_text(cand)
                    q = float(quality_score(cleaned)) if quality_score else 0.0
                    s = (0.65 * len(cleaned)) + (200.0 * q)
                    if s > best_s:
                        best_s = s
                        best = cand
                raw_text = best

            page_text = _clean_pdf_page_text(raw_text)
            if not page_text:
                continue
            if _is_probable_toc(page_text):
                continue
            full_parts.append(page_text)
            for ch in _chunk_text(page_text):
                pages_out.append({"text": ch, "meta": {"page": idx + 1}})

    full_text = _normalize_pipeline_text("\n\n".join(full_parts))
    return full_text, pages_out


def _pick_best_pdf_extraction(
    candidates: List[Tuple[str, str, List[Dict[str, Any]]]],
    *,
    total_pages: int | None = None,
) -> Optional[Tuple[str, List[Dict[str, Any]], Dict[str, Any]]]:
    """Pick best extraction with completeness-aware scoring."""
    if not candidates:
        return None

    try:
        from app.services.text_quality import quality_score
    except Exception:
        # If quality module is unavailable for some reason, fall back to longest text.
        quality_score = None  # type: ignore

    scored: List[Dict[str, Any]] = []
    for name, full_text, chunks in candidates:
        if quality_score:
            s = float(quality_score(full_text))
        else:
            s = 0.0
        scored.append(
            {
                "name": name,
                "quality_score": s,
                "char_len": len(full_text or ""),
                "full_text": full_text,
                "chunks": chunks,
                "chunk_count": len(chunks or []),
                "page_coverage": _candidate_page_coverage(chunks, total_pages=total_pages),
            }
        )

    best_cov = max((float(x["page_coverage"]) for x in scored), default=0.0)
    best_len = max((int(x["char_len"]) for x in scored), default=1)
    coverage_gate = float(getattr(settings, "PDF_EXTRACT_MIN_COVERAGE_RATIO", 0.83) or 0.83)
    eligible = [x for x in scored if best_cov <= 0 or float(x["page_coverage"]) >= (best_cov * coverage_gate)]
    if not eligible:
        eligible = scored

    for x in eligible:
        q = float(x["quality_score"])
        cov_ratio = float(x["page_coverage"]) / max(best_cov, 1e-9) if best_cov > 0 else 1.0
        len_ratio = float(x["char_len"]) / max(1.0, float(best_len))
        x["completeness_weighted_score"] = (0.38 * q) + (0.47 * cov_ratio) + (0.15 * len_ratio)

    eligible.sort(key=lambda x: (x["completeness_weighted_score"], x["quality_score"], x["char_len"]), reverse=True)
    chosen = eligible[0]
    report = {
        "chosen_extractor": chosen["name"],
        "candidates": [
            {
                "name": x["name"],
                "quality_score": round(float(x["quality_score"]), 4),
                "char_len": int(x["char_len"]),
                "page_coverage": round(float(x["page_coverage"]), 4),
                "chunk_count": int(x["chunk_count"]),
                "completeness_weighted_score": round(float(x.get("completeness_weighted_score", 0.0)), 4),
            }
            for x in scored
        ],
    }
    return chosen["full_text"], chosen["chunks"], report


def _extract_text_pdf_with_report(data: bytes) -> Tuple[str, List[Dict[str, Any]], Dict[str, Any]]:
    """Extract PDF text with multiple strategies and pick best.

    Priority per user request:
      1) pdfplumber (best for many Vietnamese textbooks)
      2) PyMuPDF (words reconstruction) + PyMuPDF text
      3) pypdf (last resort)
      4) OCR (only when text-layer quality is too low)
    """
    candidates: List[Tuple[str, str, List[Dict[str, Any]]]] = []
    total_pages = None
    try:
        import fitz

        doc = fitz.open(stream=data, filetype="pdf")
        total_pages = len(doc)
    except Exception:
        total_pages = None

    # 1) pdfplumber (preferred)
    try:
        full_text, chunks = _extract_text_pdf_pdfplumber(data)
        if full_text.strip():
            candidates.append(("pdfplumber", full_text, chunks))
    except Exception:
        pass

    # 2) PyMuPDF word reconstruction (best for missing spaces)
    try:
        full_text, chunks = _extract_text_pdf_pymupdf_words(data)
        if full_text.strip():
            candidates.append(("pymupdf_words", full_text, chunks))
    except Exception:
        pass

    # 3) PyMuPDF text
    try:
        full_text, chunks = _extract_text_pdf_pymupdf(data)
        if full_text.strip():
            candidates.append(("pymupdf", full_text, chunks))
    except Exception:
        pass

    # 4) pypdf fallback (only if nothing else worked)
    if not candidates:
        try:
            full_text, chunks = _extract_text_pdf_pypdf(data)
            if full_text.strip():
                candidates.append(("pypdf", full_text, chunks))
        except Exception:
            pass

    picked = _pick_best_pdf_extraction(candidates, total_pages=total_pages)
    if not picked:
        # last resort OCR
        ocr_text, ocr_chunks = _extract_text_pdf_ocr(data)
        if ocr_text.strip():
            return ocr_text, ocr_chunks, {
                "chosen_extractor": "ocr",
                "ocr_used": True,
                "candidates": [{"name": "ocr", "quality_score": None, "char_len": len(ocr_text), "page_coverage": _candidate_page_coverage(ocr_chunks, total_pages=total_pages), "chunk_count": len(ocr_chunks)}],
            }
        return "", [], {"chosen_extractor": None, "ocr_used": False, "candidates": []}

    best_text, best_chunks, report = picked

    # OCR trigger when extracted text is likely garbled or too sparse
    try:
        from app.services.text_quality import quality_score

        q = float(quality_score(best_text))
    except Exception:
        q = 1.0

    # Force OCR if we still detect strong spaced-letter artifacts in the sample.
    # Some PDFs have huge amounts of text, so an overall quality score can hide this issue.
    sample = (best_text or '')[:8000]
    spaced_seq_rx = re.compile(r"\b(?:[A-Za-zÀ-ỹà-ỹ]\s+){3,}[A-Za-zÀ-ỹà-ỹ]\b")
    tokens = re.findall(r"[A-Za-zÀ-ỹà-ỹ0-9]+", sample)
    one_ratio = 0.0
    if tokens:
        one_ratio = sum(1 for w in tokens if len(w) == 1 and w.isalpha()) / max(1, len(tokens))
    force_ocr = bool(spaced_seq_rx.search(sample)) or one_ratio > 0.28

    trigger = float(getattr(settings, "PDF_OCR_TRIGGER_MIN_QUALITY_SCORE", 0.22) or 0.22)
    if q < trigger or force_ocr:

        ocr_text, ocr_chunks = _extract_text_pdf_ocr(data)
        if ocr_text.strip():
            cand2 = candidates + [("ocr", ocr_text, ocr_chunks)]
            picked2 = _pick_best_pdf_extraction(cand2, total_pages=total_pages)
            if picked2:
                txt2, chunks2, report2 = picked2
                report2["ocr_used"] = bool(report2.get("chosen_extractor") == "ocr")
                return txt2, chunks2, report2

    report["ocr_used"] = False
    return best_text, best_chunks, report


def _extract_text_pdf(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    text, chunks, _ = _extract_text_pdf_with_report(data)
    return text, chunks


def _extract_text_docx(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    try:
        from docx import Document as DocxDocument
    except Exception as e:
        raise RuntimeError("Missing dependency python-docx for DOCX extraction") from e

    doc = DocxDocument(io.BytesIO(data))
    text = "\n".join([p.text for p in doc.paragraphs if p.text and p.text.strip()])
    text = _sanitize_text(text)
    chunks = [{"text": ch, "meta": {}} for ch in _chunk_text(text)]
    return text, chunks


def _extract_text_pptx(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    try:
        from pptx import Presentation
    except Exception as e:
        raise RuntimeError("Missing dependency python-pptx for PPTX extraction") from e

    prs = Presentation(io.BytesIO(data))
    slide_texts: List[str] = []
    chunks: List[Dict[str, Any]] = []

    for s_idx, slide in enumerate(prs.slides):
        lines: List[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                t = (shape.text or "").strip()
                if t:
                    lines.append(_sanitize_text(t))
        slide_text = "\n".join(lines).strip()
        if not slide_text:
            continue
        slide_texts.append(slide_text)
        for ch in _chunk_text(slide_text):
            chunks.append({"text": ch, "meta": {"slide": s_idx + 1}})

    return _normalize_pipeline_text("\n\n".join(slide_texts)), chunks


def _extract_text_fallback(data: bytes) -> Tuple[str, List[Dict[str, Any]]]:
    # Try utf-8; if binary, ignore errors.
    text = _sanitize_text(data.decode("utf-8", errors="ignore"))
    chunks = [{"text": ch, "meta": {}} for ch in _chunk_text(text)]
    return text, chunks


async def extract_and_chunk(file: UploadFile) -> Tuple[str, List[Dict[str, Any]]]:
    """Extract text from UploadFile and return (full_text, chunks[{text,meta}])."""
    data = await file.read()

    ctype = (file.content_type or "").lower()
    fname = (file.filename or "").lower()

    if ctype == "application/pdf" or fname.endswith(".pdf"):
        text, chunks, _ = _extract_text_pdf_with_report(data)
        return text, chunks
    if ctype in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    } or fname.endswith(".docx"):
        return _extract_text_docx(data)
    if ctype in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    } or fname.endswith(".pptx"):
        return _extract_text_pptx(data)

    return _extract_text_fallback(data)


async def extract_and_chunk_with_report(file: UploadFile) -> Tuple[str, List[Dict[str, Any]], Dict[str, Any] | None]:
    """Extract text/chunks and include additive extraction report for PDFs."""
    data = await file.read()
    ctype = (file.content_type or "").lower()
    fname = (file.filename or "").lower()
    if ctype == "application/pdf" or fname.endswith(".pdf"):
        text, chunks, report = _extract_text_pdf_with_report(data)
        return text, chunks, report

    if ctype in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    } or fname.endswith(".docx"):
        text, chunks = _extract_text_docx(data)
        return text, chunks, None
    if ctype in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    } or fname.endswith(".pptx"):
        text, chunks = _extract_text_pptx(data)
        return text, chunks, None
    text, chunks = _extract_text_fallback(data)
    return text, chunks, None
